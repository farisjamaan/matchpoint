"""
services/ingestion.py

Handles the two-phase data ingestion pipeline:
  1. parse_pptx        — extract plain text from a .pptx file
  2. parse_text_file   — extract plain text from a slide-marker .txt file
  3. extract_metadata_with_llm  — call Groq to pull structured profile data
  4. ingest_all        — orchestrate both phases across .pptx and .txt files

Dependencies (SQLiteManager, Groq client, model names) are injected by the
caller; nothing is instantiated globally inside this module.
"""

from __future__ import annotations

import json
import re
from pathlib import Path

from groq import Groq
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from src.core.logger import get_logger
from src.db.sqlite_mgr import SQLiteManager

logger = get_logger(__name__)


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


def parse_pptx(filepath: Path) -> str:
    """
    Extract all text from a .pptx file as plain text.

    Each shape's lines are grouped into a block; blocks are separated by
    double newlines so downstream chunking (split on "\\n\\n") produces
    one chunk per shape rather than one giant chunk per candidate.
    Raises PackageNotFoundError if the file is not a valid .pptx.
    """
    prs = Presentation(str(filepath))
    shape_blocks: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            lines = [line.strip() for line in shape.text.strip().split("\n") if line.strip()]
            if lines:
                shape_blocks.append("\n".join(lines))

    return "\n\n".join(shape_blocks)


# ---------------------------------------------------------------------------
# Plain-text (slide-marker format) parser — supports multi-person files
# ---------------------------------------------------------------------------

# A profile-header slide contains an email address or an international-style
# phone number, which reliably marks the start of a new person's resume.
_EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}")
_PHONE_RE = re.compile(r"\+?\d[\d\s\-\(\)]{6,}\d")
_PAGE_MARKER_RE = re.compile(r"^Page\s+\d+$")


def _slides_from_text(raw: str) -> list[list[str]]:
    """
    Split raw text into individual slides.
    Each slide is returned as a list of non-empty content lines.
    Slide-marker fences (=====) and 'Slide N' header lines are discarded.
    """
    _SEPARATOR = "======================"
    slides: list[list[str]] = []
    current: list[str] = []

    for line in raw.splitlines():
        stripped = line.strip()
        if stripped == _SEPARATOR:
            continue
        parts = stripped.split()
        if len(parts) == 2 and parts[0] == "Slide" and parts[1].isdigit():
            if current:
                slides.append(current)
                current = []
            continue
        if stripped:
            current.append(stripped)

    if current:
        slides.append(current)

    return slides


def _is_profile_start(slide_lines: list[str]) -> bool:
    """
    Return True when a slide looks like the opening profile slide for a new
    person — i.e. it contains an email address or phone number.
    """
    text = " ".join(slide_lines)
    return bool(_EMAIL_RE.search(text) or _PHONE_RE.search(text))


def parse_text_file(filepath: Path) -> list[tuple[str, str]]:
    """
    Parse a plain-text CV file (slide-marker format) that may contain one or
    more people's resumes.

    Returns a list of (synthetic_filename, content) pairs — one per person.
    A new person boundary is detected when a slide contains an email address
    or phone number (the profile-header slide pattern).

    The synthetic filename is built from the original stem so the SQLite
    upsert key remains unique per person:
        pablo_perez_cv.txt  →  pablo_perez_cv_person1.txt, _person2.txt, …
    """
    stem = filepath.stem
    ext = filepath.suffix
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    slides = _slides_from_text(raw)

    # Group slides into per-person buckets.
    # A new bucket starts whenever a profile-header slide is detected AND
    # there is already content in the current bucket.
    person_groups: list[list[list[str]]] = []
    current_group: list[list[str]] = []

    for i, slide_lines in enumerate(slides):
        # A new person starts when:
        # 1. The slide contains an email or phone (explicit contact info), OR
        # 2. The previous slide ended with "Page N" (every candidate's last line),
        #    which means this slide is unconditionally a new person even if they
        #    have no contact info in their slide.
        prev_ended_with_page = (
            i > 0
            and slides[i - 1]
            and bool(_PAGE_MARKER_RE.match(slides[i - 1][-1]))
        )
        if (_is_profile_start(slide_lines) or prev_ended_with_page) and current_group:
            person_groups.append(current_group)
            current_group = [slide_lines]
        else:
            current_group.append(slide_lines)

    if current_group:
        person_groups.append(current_group)

    # Convert each group to a (synthetic_filename, content) tuple.
    results: list[tuple[str, str]] = []
    for i, group in enumerate(person_groups):
        content = "\n\n".join("\n".join(lines) for lines in group)
        if not content.strip():
            continue
        synthetic_name = f"{stem}_person{i + 1}{ext}"
        results.append((synthetic_name, content))

    return results


# ---------------------------------------------------------------------------
# LLM metadata extractor
# ---------------------------------------------------------------------------


def extract_metadata_with_llm(
    cv_content: str,
    groq_client: Groq,
    extraction_model: str,
    temperature: float,
) -> dict[str, str | None]:
    """
    Ask the LLM to extract {name, role, email, phone} from the first 2 000
    characters of the CV content.

    Returns a dict with those four keys; falls back to safe defaults on any
    Groq or JSON parsing error without crashing the ingestion pipeline.
    """
    _FALLBACK: dict[str, str | None] = {
        "name": "Unknown",
        "role": None,
        "email": None,
        "phone": None,
    }

    prompt = (
        "Extract core profile details. "
        "Return ONLY a valid JSON object with keys: "
        '"name", "role", "email", "phone".\n'
        f"CV Text:\n{cv_content[:2000]}"
    )

    try:
        response = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=extraction_model,
            temperature=temperature,
            response_format={"type": "json_object"},
        )
        return json.loads(response.choices[0].message.content)

    except json.JSONDecodeError as exc:
        logger.error("LLM returned invalid JSON during metadata extraction: %s", exc)
        return _FALLBACK

    except Exception as exc:  # Groq API errors, network issues, etc.
        logger.error("Groq API call failed during metadata extraction: %s", exc)
        return _FALLBACK


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------


def ingest_all(
    data_dir: Path,
    db_manager: SQLiteManager,
    groq_client: Groq,
    extraction_model: str,
    temperature: float,
) -> int:
    """
    Scan data_dir for .pptx and .txt files, parse + extract metadata for each,
    and upsert into SQLite.

    .pptx files are parsed with parse_pptx.
    .txt files are expected to use the slide-marker format and parsed with
    parse_text_file.

    Returns the count of successfully ingested files.
    Raises FileNotFoundError if data_dir does not exist.
    """
    if not data_dir.exists():
        raise FileNotFoundError(
            f"Data directory not found: {data_dir}. "
            "Create it and add .pptx or .txt files before calling /ingest."
        )

    all_files = list(data_dir.glob("*.pptx")) + list(data_dir.glob("*.txt"))
    if not all_files:
        logger.warning("No .pptx or .txt files found in %s.", data_dir)
        return 0

    logger.info("Found %d file(s) to ingest.", len(all_files))
    ingested_count = 0

    for filepath in all_files:
        logger.info("Processing: %s", filepath.name)

        # Build a list of (filename_key, content) pairs.
        # PPTX → always one person; .txt → one or more people detected by
        # the profile-header heuristic.
        try:
            if filepath.suffix.lower() == ".pptx":
                person_records: list[tuple[str, str]] = [(filepath.name, parse_pptx(filepath))]
            else:
                person_records = parse_text_file(filepath)
                logger.info(
                    "%s: detected %d person(s).", filepath.name, len(person_records)
                )
        except PackageNotFoundError:
            logger.error("Skipping %s — not a valid .pptx file.", filepath.name)
            continue
        except Exception as exc:
            logger.error("Failed to parse %s: %s", filepath.name, exc)
            continue

        for synthetic_filename, content in person_records:
            metadata = extract_metadata_with_llm(
                cv_content=content,
                groq_client=groq_client,
                extraction_model=extraction_model,
                temperature=temperature,
            )

            try:
                db_manager.upsert_candidate(
                    filename=synthetic_filename,
                    name=metadata.get("name") or "Unknown",
                    role=metadata.get("role"),
                    email=metadata.get("email"),
                    phone=metadata.get("phone"),
                    content=content,
                )
                ingested_count += 1
            except Exception as exc:
                logger.error("Failed to write %s to SQLite: %s", synthetic_filename, exc)

    logger.info(
        "Ingestion complete: %d candidate(s) stored from %d file(s).",
        ingested_count,
        len(all_files),
    )
    return ingested_count


# Backwards-compat alias for any callers using the old name
ingest_all_pptx = ingest_all
